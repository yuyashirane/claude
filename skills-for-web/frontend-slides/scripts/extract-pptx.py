#!/usr/bin/env python3
"""
extract-pptx.py — Extract slide content from a .pptx file to JSON.

Usage:
    python3 extract-pptx.py <input.pptx> [--out slides.json] [--images]

Output JSON schema:
    [
      {
        "index":  0,
        "title":  "Slide title text",
        "body":   ["Bullet one", "Bullet two"],
        "notes":  "Speaker notes text",
        "images": ["data:image/png;base64,..."]   # only with --images flag
      },
      ...
    ]

Requirements:
    pip install python-pptx Pillow

Options:
    --out PATH      Output JSON file path (default: slides.json)
    --images        Embed slide images as base64 data URIs
    --max-img-kb N  Skip images larger than N KB (default: 200)
    --pretty        Pretty-print JSON output
"""

import argparse
import base64
import io
import json
import sys
from pathlib import Path


def require_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Pt
        return Presentation
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
        sys.exit(1)


def shape_text(shape) -> list[str]:
    """Return non-empty paragraph strings from a text-bearing shape."""
    lines = []
    if not shape.has_text_frame:
        return lines
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            lines.append(text)
    return lines


def extract_title(slide) -> str:
    """Return the title placeholder text, or empty string."""
    try:
        ph = slide.shapes.title
        if ph and ph.has_text_frame:
            return ph.text_frame.text.strip()
    except Exception:
        pass
    return ""


def extract_body(slide, title_text: str) -> list[str]:
    """Return all text lines from non-title shapes."""
    lines = []
    for shape in slide.shapes:
        # Skip the title placeholder
        try:
            if shape == slide.shapes.title:
                continue
        except Exception:
            pass

        lines.extend(shape_text(shape))

    # De-duplicate while preserving order
    seen = set()
    deduped = []
    for line in lines:
        if line not in seen:
            seen.add(line)
            deduped.append(line)

    return deduped


def extract_notes(slide) -> str:
    """Return speaker notes text."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        return tf.text.strip() if tf else ""
    except Exception:
        return ""


def image_to_data_uri(image_bytes: bytes, content_type: str, max_kb: int) -> str | None:
    """Convert raw image bytes to a base64 data URI, or None if too large."""
    if len(image_bytes) > max_kb * 1024:
        return None
    b64 = base64.b64encode(image_bytes).decode("ascii")
    return f"data:{content_type};base64,{b64}"


def extract_images(slide, max_kb: int) -> list[str]:
    """Extract embedded images from a slide as base64 data URIs."""
    uris = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                img = shape.image
                uri = image_to_data_uri(img.blob, img.content_type, max_kb)
                if uri:
                    uris.append(uri)
            except Exception:
                pass
    return uris


def extract(pptx_path: Path, include_images: bool, max_img_kb: int) -> list[dict]:
    Presentation = require_pptx()
    prs = Presentation(str(pptx_path))
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        title = extract_title(slide)
        body  = extract_body(slide, title)
        notes = extract_notes(slide)

        entry: dict = {
            "index": idx,
            "title": title,
            "body":  body,
            "notes": notes,
        }

        if include_images:
            entry["images"] = extract_images(slide, max_img_kb)

        slides_data.append(entry)

    return slides_data


def main():
    parser = argparse.ArgumentParser(
        description="Extract PPTX slide content to JSON for the frontend-slides skill."
    )
    parser.add_argument("input",         type=Path,         help="Path to the .pptx file")
    parser.add_argument("--out",         type=Path,         default=Path("slides.json"),
                        help="Output JSON file path (default: slides.json)")
    parser.add_argument("--images",      action="store_true",
                        help="Embed images as base64 data URIs")
    parser.add_argument("--max-img-kb",  type=int,          default=200,
                        help="Max image size in KB to embed (default: 200)")
    parser.add_argument("--pretty",      action="store_true",
                        help="Pretty-print JSON output")

    args = parser.parse_args()

    if not args.input.exists():
        print(f"ERROR: File not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    if args.input.suffix.lower() not in (".pptx", ".ppt"):
        print(f"WARNING: Expected a .pptx file, got: {args.input.suffix}", file=sys.stderr)

    print(f"Extracting: {args.input} …")
    slides_data = extract(args.input, args.images, args.max_img_kb)

    indent = 2 if args.pretty else None
    output_json = json.dumps(slides_data, ensure_ascii=False, indent=indent)

    args.out.write_text(output_json, encoding="utf-8")

    print(f"Done. {len(slides_data)} slide(s) written to: {args.out}")
    if args.images:
        total_imgs = sum(len(s.get("images", [])) for s in slides_data)
        print(f"       {total_imgs} image(s) embedded.")


if __name__ == "__main__":
    main()
